"""
Generate presentation slides (.pptx) for the Task Management System project.
Import the output file into Google Slides via File > Import slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
C_DARK_BG    = RGBColor(0x1E, 0x1E, 0x2E)   # deep navy
C_ACCENT     = RGBColor(0x7C, 0x3A, 0xED)   # purple
C_ACCENT2    = RGBColor(0xA7, 0x8B, 0xFA)   # light purple
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xE5, 0xE7, 0xEB)
C_MID_GRAY   = RGBColor(0x9C, 0xA3, 0xAF)
C_CARD       = RGBColor(0x2D, 0x2D, 0x44)   # card background
C_GREEN      = RGBColor(0x10, 0xB9, 0x81)
C_YELLOW     = RGBColor(0xF5, 0x9E, 0x0B)
C_RED        = RGBColor(0xEF, 0x44, 0x44)
C_BLUE       = RGBColor(0x38, 0xBD, 0xF8)
C_ORANGE     = RGBColor(0xF9, 0x73, 0x16)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]   # totally blank


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width if line_width else Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=Pt(18), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def bg(slide, color=C_DARK_BG):
    """Fill slide background."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=color)


def accent_bar(slide, height=Inches(0.08), color=C_ACCENT):
    """Top accent bar."""
    add_rect(slide, 0, 0, SLIDE_W, height, fill=color)


def section_header(slide, title, subtitle=None):
    """Standard section header block."""
    add_rect(slide, 0, Inches(2.8), SLIDE_W, Inches(2.0), fill=C_CARD)
    add_text(slide, title,
             Inches(1.5), Inches(3.0), Inches(10), Inches(0.9),
             font_size=Pt(44), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(1.5), Inches(4.0), Inches(10), Inches(0.5),
                 font_size=Pt(20), color=C_ACCENT2, align=PP_ALIGN.CENTER)


def slide_title(slide, title, y=Inches(0.25)):
    add_text(slide, title,
             Inches(0.5), y, Inches(12.3), Inches(0.6),
             font_size=Pt(28), bold=True, color=C_WHITE)
    # underline rule
    add_rect(slide, Inches(0.5), y + Inches(0.65), Inches(12.3), Inches(0.04),
             fill=C_ACCENT)


def bullet_card(slide, x, y, w, h, title, bullets, title_color=C_ACCENT2):
    add_rect(slide, x, y, w, h, fill=C_CARD, line_color=C_ACCENT, line_width=Pt(1))
    add_text(slide, title,
             x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.4),
             font_size=Pt(14), bold=True, color=title_color)
    body = "\n".join(f"• {b}" for b in bullets)
    add_text(slide, body,
             x + Inches(0.15), y + Inches(0.5), w - Inches(0.3), h - Inches(0.6),
             font_size=Pt(12), color=C_LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
# gradient-style horizontal band
add_rect(s, 0, Inches(1.8), SLIDE_W, Inches(3.9), fill=C_CARD)
add_rect(s, 0, Inches(1.8), Inches(0.25), Inches(3.9), fill=C_ACCENT)

add_text(s, "TASK MANAGEMENT SYSTEM",
         Inches(0.6), Inches(2.0), Inches(12), Inches(0.9),
         font_size=Pt(46), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s, "AI-Assisted Development with Claude Code",
         Inches(0.6), Inches(3.0), Inches(12), Inches(0.6),
         font_size=Pt(26), color=C_ACCENT2, align=PP_ALIGN.CENTER)
add_text(s, "A Research Experiment — Anthropic Claude Code Pro Plan",
         Inches(0.6), Inches(3.7), Inches(12), Inches(0.45),
         font_size=Pt(18), color=C_MID_GRAY, align=PP_ALIGN.CENTER, italic=True)

# Tag chips
chips = [("Next.js 16", C_ACCENT), ("PostgreSQL", C_BLUE), ("Prisma ORM", C_GREEN),
         ("TypeScript", C_YELLOW), ("Tailwind CSS", C_ORANGE)]
cx = Inches(2.1)
for label, col in chips:
    add_rect(s, cx, Inches(4.6), Inches(1.7), Inches(0.38), fill=col)
    add_text(s, label, cx + Inches(0.08), Inches(4.62), Inches(1.55), Inches(0.34),
             font_size=Pt(13), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    cx += Inches(1.85)

add_text(s, "March 2026  |  Version 1",
         Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
         font_size=Pt(14), color=C_MID_GRAY, align=PP_ALIGN.CENTER)
add_rect(s, 0, Inches(7.3), SLIDE_W, Inches(0.2), fill=C_ACCENT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "Agenda")

items = [
    ("01", "Project Background",         "What the experiment is and why it was done"),
    ("02", "Development Timeline",        "5 phases from scaffold to production"),
    ("03", "Feature Evolution",           "How the system expanded over time"),
    ("04", "Bug → Fix Mapping",           "17 bugs documented with root causes and fixes"),
    ("05", "Prompt → Result Examples",    "How AI interaction shaped the code"),
    ("06", "Results & Observations",      "What we learned about AI-assisted development"),
    ("07", "Conclusion & Next Steps",     "Version 2 roadmap and final thoughts"),
]

cols = 2
rows_per_col = 4
x_starts = [Inches(0.5), Inches(6.9)]
for i, (num, title, desc) in enumerate(items):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    x = x_starts[col]
    y = Inches(1.3) + row * Inches(1.4)
    add_rect(s, x, y, Inches(5.9), Inches(1.2), fill=C_CARD, line_color=C_ACCENT, line_width=Pt(1))
    add_rect(s, x, y, Inches(0.5), Inches(1.2), fill=C_ACCENT)
    add_text(s, num, x + Inches(0.02), y + Inches(0.35), Inches(0.46), Inches(0.5),
             font_size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, x + Inches(0.6), y + Inches(0.1), Inches(5.1), Inches(0.4),
             font_size=Pt(15), bold=True, color=C_ACCENT2)
    add_text(s, desc, x + Inches(0.6), y + Inches(0.5), Inches(5.1), Inches(0.6),
             font_size=Pt(12), color=C_LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Project Background
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "01 — Project Background")

add_text(s, "Why This Experiment?",
         Inches(0.5), Inches(1.2), Inches(12), Inches(0.45),
         font_size=Pt(20), bold=True, color=C_ACCENT2)
add_text(s,
    "Instead of analysing Claude Code theoretically, this experiment built a real working system to "
    "evaluate whether AI can follow proper coding practices, increase development efficiency, and "
    "produce production-level software — not just simple demos.",
    Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.9),
    font_size=Pt(14), color=C_LIGHT_GRAY)

goals = [
    ("🎯 Research Goal",   ["Understand real-world capability of AI coding tools",
                             "Evaluate code quality, architecture, and best practices",
                             "Measure development speed vs traditional approach"]),
    ("🛠 Tool Used",       ["Claude Code by Anthropic (Pro Plan)",
                             "Natural language prompts → working code",
                             "4-day development session: March 12–16, 2026"]),
    ("📐 Scope",           ["Full-stack web application — not a toy project",
                             "Production features: auth, CRUD, real-time UI, email",
                             "Version 1 complete; Version 2 planned"]),
    ("📊 Evaluation",      ["Does the AI follow software engineering best practices?",
                             "Is the output closer to a demo or production software?",
                             "What human oversight is still required?"]),
]
xp = [Inches(0.5), Inches(3.65), Inches(6.8), Inches(10.0)]
for i, (title, bullets) in enumerate(goals):
    bullet_card(s, xp[i], Inches(2.8), Inches(3.0), Inches(2.0), title, bullets)

add_rect(s, 0, Inches(6.9), SLIDE_W, Inches(0.6), fill=C_CARD)
add_text(s,
    "\"The experiment focused on building a working system to evaluate whether AI can assist in the full software development lifecycle.\"",
    Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.5),
    font_size=Pt(13), color=C_ACCENT2, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — System Overview
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "System Overview — What Was Built")

features = [
    ("🔐 Authentication",  ["Email/password + magic link login","JWT sessions (30-day)","Email invite flow","Role-based access (admin/member)"]),
    ("📋 Task Management", ["Kanban board with drag-and-drop","List view with sort & filter","Multi-assignee (up to 5)","Past-date due date validation"]),
    ("💬 Collaboration",   ["Threaded task comments","Board discussion channel","Emoji reactions","@mention + email notifications"]),
    ("👤 User Management", ["Admin CRUD for users","Invite/revoke flow","Deactivate/reactivate users","Profile photos (WebP 256×256)"]),
    ("🎨 UI/UX",           ["JIRA-style two-column task modal","Sonner toast notifications","Confirm dialogs","Pagination (10/page)"]),
    ("🗄 Tech Stack",       ["Next.js 16 + React 19 + TypeScript","PostgreSQL + Prisma 7 ORM","Tailwind CSS v4","@dnd-kit drag-and-drop"]),
]
xpos = [Inches(0.4), Inches(4.65), Inches(8.9)]
for i, (title, bullets) in enumerate(features):
    col = i % 3
    row = i // 3
    x = xpos[col]
    y = Inches(1.3) + row * Inches(2.4)
    bullet_card(s, x, y, Inches(4.0), Inches(2.2), title, bullets)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Development Timeline (table style)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "02 — Development Timeline")

phases = [
    ("Phase 1", "Foundation",       "March 12",  "Project scaffold, Auth API, User Management API, Task API, DB fixes",         C_ACCENT),
    ("Phase 2", "Core Frontend",    "March 13",  "Jira-inspired UI, Kanban board, Invite flow, Magic link, Multi-assignee",      C_BLUE),
    ("Phase 3", "Enrichment",       "March 13",  "Profile photo upload, Avatar display, Emoji reactions, Optimistic UI",         C_GREEN),
    ("Phase 4", "Advanced Features","March 16",  "Pagination, Activity feed, Pin comments, @Mentions, JIRA modal, Toasts",       C_YELLOW),
    ("Phase 5", "Validation",       "March 16",  "Past due date prevention, Admin self-deactivation guard, Bug hardening",       C_ORANGE),
]

# Header row
hx = [Inches(0.4), Inches(1.5), Inches(3.0), Inches(4.2), Inches(6.3)]
hw = [Inches(1.0), Inches(1.4), Inches(1.1), Inches(2.0), Inches(6.5)]
headers = ["Phase", "Name", "Date", "Focus", "Key Deliverables"]
add_rect(s, Inches(0.4), Inches(1.3), Inches(12.4), Inches(0.42), fill=C_ACCENT)
for j, h in enumerate(headers):
    add_text(s, h, hx[j], Inches(1.32), hw[j], Inches(0.38),
             font_size=Pt(12), bold=True, color=C_WHITE)

for i, (phase, name, date, focus, col) in enumerate(phases):
    ry = Inches(1.72) + i * Inches(0.96)
    row_bg = C_CARD if i % 2 == 0 else RGBColor(0x25, 0x25, 0x3A)
    add_rect(s, Inches(0.4), ry, Inches(12.4), Inches(0.9), fill=row_bg)
    add_rect(s, Inches(0.4), ry, Inches(0.18), Inches(0.9), fill=col)
    data = [phase, name, date, focus]
    for j, txt in enumerate(data):
        add_text(s, txt, hx[j] + Inches(0.05), ry + Inches(0.05),
                 hw[j] - Inches(0.1), Inches(0.8),
                 font_size=Pt(11), color=C_LIGHT_GRAY,
                 bold=(j == 0))

# Stats bar
add_rect(s, 0, Inches(6.95), SLIDE_W, Inches(0.55), fill=C_CARD)
stats = [("46", "Git Commits"), ("27", "Development Steps"), ("17", "Bugs Fixed"), ("4", "Days")]
sx = Inches(1.2)
for val, lbl in stats:
    add_text(s, val, sx, Inches(6.98), Inches(1.2), Inches(0.28),
             font_size=Pt(20), bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)
    add_text(s, lbl, sx, Inches(7.24), Inches(1.2), Inches(0.22),
             font_size=Pt(10), color=C_MID_GRAY, align=PP_ALIGN.CENTER)
    sx += Inches(2.6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Feature Evolution Diagram
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "03 — Feature Evolution Diagram")

layers = [
    ("Foundation Layer",
     ["Project Scaffold", "JWT Auth", "User CRUD API", "Task CRUD API", "Prisma DB"],
     C_ACCENT, Inches(1.25)),
    ("Core Frontend",
     ["Kanban Board", "List View", "Task Modal", "Comments", "User Mgmt Table"],
     C_BLUE, Inches(2.35)),
    ("Enrichment",
     ["Profile Photos", "Emoji Reactions", "Multi-Assignee", "Board Discussion", "Optimistic UI"],
     C_GREEN, Inches(3.45)),
    ("Advanced Features",
     ["@Mentions + Email", "JIRA Modal", "Pin Comments", "Toast + Dialogs", "Pagination"],
     C_YELLOW, Inches(4.55)),
    ("Version 2 (Planned)",
     ["SSO Authentication", "Mobile Responsive", "Custom Task Fields", "Analytics Charts", "Sub-tasks"],
     RGBColor(0x6B, 0x72, 0x80), Inches(5.65)),
]

for i, (layer_name, items, col, y) in enumerate(layers):
    is_planned = "Planned" in layer_name
    line_col = RGBColor(0x6B, 0x72, 0x80) if is_planned else col
    fill_col = RGBColor(0x28, 0x28, 0x3A) if is_planned else C_CARD
    add_rect(s, Inches(0.35), y, Inches(12.6), Inches(0.9),
             fill=fill_col, line_color=line_col, line_width=Pt(1))
    add_rect(s, Inches(0.35), y, Inches(0.18), Inches(0.9), fill=col)
    add_text(s, layer_name,
             Inches(0.65), y + Inches(0.08), Inches(1.9), Inches(0.4),
             font_size=Pt(11), bold=True,
             color=C_MID_GRAY if is_planned else col)
    bx = Inches(2.7)
    for item in items:
        item_col = RGBColor(0x4B, 0x55, 0x63) if is_planned else col
        bg_col   = RGBColor(0x2D, 0x2D, 0x3A) if is_planned else C_DARK_BG
        add_rect(s, bx, y + Inches(0.1), Inches(1.95), Inches(0.7),
                 fill=bg_col, line_color=item_col, line_width=Pt(1))
        add_text(s, item, bx + Inches(0.08), y + Inches(0.2), Inches(1.8), Inches(0.5),
                 font_size=Pt(10), color=C_MID_GRAY if is_planned else C_LIGHT_GRAY,
                 align=PP_ALIGN.CENTER, italic=is_planned)
        bx += Inches(2.05)

# Arrow
add_text(s, "▲ V1 Shipped", Inches(11.8), Inches(4.8), Inches(1.4), Inches(0.35),
         font_size=Pt(11), bold=True, color=C_GREEN)
add_text(s, "▲ V2 Planned", Inches(11.8), Inches(5.6), Inches(1.4), Inches(0.35),
         font_size=Pt(11), bold=True, color=C_MID_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Bug → Fix Mapping (Part 1)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "04 — Bug → Fix Mapping Table  (Part 1 of 2)")

bugs_p1 = [
    ("DB-01", "Seed crashes — DATABASE_URL undefined",   "dotenv loaded after PrismaClient import",          "Move dotenv/config to first line of seed.ts"),
    ("DB-02", "Prisma can't find datasource URL",        "Prisma v7 requires explicit url in schema",         "Add url = env(\"DATABASE_URL\") to datasource block"),
    ("DB-03", "Prisma runtime driver error",             "Prisma v7 requires explicit pg adapter",            "Pass @prisma/adapter-pg to PrismaClient constructor"),
    ("API-01","400 creating unassigned task",            "Validation rejected null assignedUserId",           "Allow null in task creation validation"),
    ("UI-01", "To Do column border invisible",           "Background and border CSS classes conflicted",      "Split into separate bg/border colour classes"),
    ("API-02","Runtime crash on empty API response",     ".json() called unconditionally",                    "Guard with content-type check before parsing"),
    ("RX-01", "Duplicate emoji reactions submitted",     "No in-flight guard on reaction requests",          "Added in-flight Set to block concurrent requests"),
    ("UI-02", "Empty board on fetch failure",            "Error silently swallowed in catch block",           "Set error state; render banner with Retry button"),
    ("RX-02", "Reaction button disabled state inverted", "Boolean expression negated incorrectly",            "Corrected the disabled prop logic"),
]

col_x = [Inches(0.4), Inches(1.15), Inches(3.5), Inches(7.3)]
col_w = [Inches(0.7), Inches(2.25), Inches(3.7), Inches(5.3)]
headers = ["ID", "Bug", "Root Cause", "Fix Applied"]

add_rect(s, Inches(0.4), Inches(1.3), Inches(12.4), Inches(0.38), fill=C_ACCENT)
for j, h in enumerate(headers):
    add_text(s, h, col_x[j], Inches(1.32), col_w[j], Inches(0.34),
             font_size=Pt(11), bold=True, color=C_WHITE)

for i, (bid, bug, cause, fix) in enumerate(bugs_p1):
    ry = Inches(1.68) + i * Inches(0.6)
    rbg = C_CARD if i % 2 == 0 else RGBColor(0x25, 0x25, 0x3A)
    add_rect(s, Inches(0.4), ry, Inches(12.4), Inches(0.56), fill=rbg)
    id_col = C_RED if bid.startswith("DB") else (C_YELLOW if bid.startswith("API") else C_BLUE)
    add_text(s, bid, col_x[0], ry + Inches(0.05), col_w[0], Inches(0.46),
             font_size=Pt(9), bold=True, color=id_col)
    add_text(s, bug,   col_x[1], ry + Inches(0.05), col_w[1], Inches(0.46),
             font_size=Pt(9), color=C_LIGHT_GRAY)
    add_text(s, cause, col_x[2], ry + Inches(0.05), col_w[2], Inches(0.46),
             font_size=Pt(9), color=C_MID_GRAY, italic=True)
    add_text(s, fix,   col_x[3], ry + Inches(0.05), col_w[3], Inches(0.46),
             font_size=Pt(9), color=C_GREEN)

add_text(s, "9 of 17 bugs shown — continued on next slide",
         Inches(0.4), Inches(7.2), Inches(12), Inches(0.3),
         font_size=Pt(11), color=C_MID_GRAY, italic=True, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Bug → Fix Mapping (Part 2)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "04 — Bug → Fix Mapping Table  (Part 2 of 2)")

bugs_p2 = [
    ("BD-01", "GET 500 after pin schema change",         "Stale cached Prisma client missing new fields",    "Added try-catch fallback query omitting pinned fields"),
    ("BD-02", "Input cleared even when POST fails",      "Clear/re-fetch not gated on response.ok",          "Added if (!response.ok) return guard"),
    ("MN-01", "Mention dropdown invisible in board",     "Clipped by overflow:hidden on parent container",   "Rendered dropdown via createPortal with position:fixed"),
    ("MN-02", "'undefined mentioned you' in email",      "Middleware not forwarding x-user-name header",     "Added x-user-name injection in middleware"),
    ("DB-04", "P2022 column not found on POST",          "pinned/pinnedAt in schema but migration not run",  "Generated and applied Prisma migration for columns"),
    ("BD-03", "Pin state resets after task interaction", "Promise.all failure wiped board comment state",    "Decoupled board and task comment fetches independently"),
    ("DT-01", "Past dates accepted in task creation",    "No min-date validation in UI or API",              "Added min=today to input + server-side 422 validation"),
    ("UM-01", "Admin could deactivate own account",      "No self-edit guard on status field",               "Added guard comparing target ID to caller ID → 403"),
]

add_rect(s, Inches(0.4), Inches(1.3), Inches(12.4), Inches(0.38), fill=C_ACCENT)
for j, h in enumerate(headers):
    add_text(s, h, col_x[j], Inches(1.32), col_w[j], Inches(0.34),
             font_size=Pt(11), bold=True, color=C_WHITE)

for i, (bid, bug, cause, fix) in enumerate(bugs_p2):
    ry = Inches(1.68) + i * Inches(0.66)
    rbg = C_CARD if i % 2 == 0 else RGBColor(0x25, 0x25, 0x3A)
    add_rect(s, Inches(0.4), ry, Inches(12.4), Inches(0.62), fill=rbg)
    id_col = C_GREEN if bid.startswith("BD") else (C_ORANGE if bid.startswith("MN") else C_ACCENT2)
    add_text(s, bid, col_x[0], ry + Inches(0.05), col_w[0], Inches(0.52),
             font_size=Pt(9), bold=True, color=id_col)
    add_text(s, bug,   col_x[1], ry + Inches(0.05), col_w[1], Inches(0.52),
             font_size=Pt(9), color=C_LIGHT_GRAY)
    add_text(s, cause, col_x[2], ry + Inches(0.05), col_w[2], Inches(0.52),
             font_size=Pt(9), color=C_MID_GRAY, italic=True)
    add_text(s, fix,   col_x[3], ry + Inches(0.05), col_w[3], Inches(0.52),
             font_size=Pt(9), color=C_GREEN)

# Summary pill
add_rect(s, Inches(0.4), Inches(7.0), Inches(12.4), Inches(0.42), fill=C_CARD)
add_text(s, "17 bugs total  ·  All resolved  ·  Each fix verified before next feature was added",
         Inches(0.6), Inches(7.05), Inches(12), Inches(0.32),
         font_size=Pt(12), color=C_ACCENT2, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Prompt → Result Examples
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "05 — Prompt → Result Examples")

examples = [
    {
        "step": "Step 1 — Scaffold",
        "prompt": "Set up a full-stack RIO Task web application using Next.js, PostgreSQL, and Prisma with TypeScript and Tailwind CSS.",
        "result": "Generated complete project structure: package.json, Prisma schema with 6 models, environment config, seed file with default admin account, and TypeScript + Tailwind setup.",
        "col": C_ACCENT,
    },
    {
        "step": "Step 2 — Feature",
        "prompt": "Build a Jira-inspired UI with a Kanban board, drag-and-drop, task CRUD modal, comments section, and user management table.",
        "result": "Produced 5 full React components (TaskBoard, TaskModal, CommentSection, UserManagementTable, Sidebar) with @dnd-kit integration, optimistic updates, and role-based UI.",
        "col": C_BLUE,
    },
    {
        "step": "Step 3 — Bug Fix",
        "prompt": "Fix the @mention dropdown so it works in the board discussion area. The dropdown is being clipped by overflow:hidden.",
        "result": "Identified root cause (overflow container), extracted MentionTextarea component, re-implemented dropdown as a portal at document.body with position:fixed using getBoundingClientRect().",
        "col": C_GREEN,
    },
    {
        "step": "Step 4 — Refinement",
        "prompt": "Redesign the task modal in JIRA style with a two-column layout. Left: title, description, activity. Right: status, priority, due date, assignees.",
        "result": "Rebuilt TaskModal in edit mode as a wide two-column layout, preserved compact single-column for create mode, added live status badge and priority dot in header.",
        "col": C_YELLOW,
    },
]

for i, ex in enumerate(examples):
    col = 0 if i < 2 else 1
    row = i % 2
    x = Inches(0.4) if col == 0 else Inches(6.7)
    y = Inches(1.3) + row * Inches(2.75)
    add_rect(s, x, y, Inches(6.1), Inches(2.55), fill=C_CARD,
             line_color=ex["col"], line_width=Pt(1))
    add_rect(s, x, y, Inches(6.1), Inches(0.32), fill=ex["col"])
    add_text(s, ex["step"], x + Inches(0.12), y + Inches(0.04),
             Inches(5.8), Inches(0.28), font_size=Pt(11), bold=True, color=C_WHITE)
    add_text(s, "PROMPT:", x + Inches(0.12), y + Inches(0.38),
             Inches(1.0), Inches(0.22), font_size=Pt(9), bold=True, color=ex["col"])
    add_text(s, ex["prompt"], x + Inches(0.12), y + Inches(0.58),
             Inches(5.85), Inches(0.78), font_size=Pt(10), color=C_LIGHT_GRAY, italic=True)
    add_text(s, "RESULT:", x + Inches(0.12), y + Inches(1.38),
             Inches(1.0), Inches(0.22), font_size=Pt(9), bold=True, color=C_GREEN)
    add_text(s, ex["result"], x + Inches(0.12), y + Inches(1.58),
             Inches(5.85), Inches(0.88), font_size=Pt(10), color=C_LIGHT_GRAY)

add_text(s, "Each feature followed an iterative prompt → test → refine cycle.",
         Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.3),
         font_size=Pt(12), color=C_MID_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Results & Observations
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "06 — Results & Observations")

add_text(s, "What We Found",
         Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
         font_size=Pt(20), bold=True, color=C_ACCENT2)

strengths = [
    "Rapidly scaffolded full-stack applications from a single prompt",
    "Consistently produced working, runnable code on the first attempt",
    "Efficiently debugged issues when given clear error descriptions",
    "Implemented complex features (drag-and-drop, portals, optimistic UI) correctly",
    "Maintained code consistency across 46 commits and 4 development phases",
    "Handled iterative refinement well — each follow-up prompt improved the output",
]
limitations = [
    "Prioritised making code work over enforcing architectural best practices",
    "Required human direction to decide when and what to build next",
    "Occasionally introduced technical debt that needed follow-up correction",
    "Could not self-detect missing database migrations without being told",
    "No automated testing was written — all verification was manual",
    "Architecture decisions still required developer review and guidance",
]

# Strengths card
add_rect(s, Inches(0.4), Inches(1.65), Inches(6.05), Inches(4.8), fill=C_CARD,
         line_color=C_GREEN, line_width=Pt(1.5))
add_rect(s, Inches(0.4), Inches(1.65), Inches(6.05), Inches(0.38), fill=C_GREEN)
add_text(s, "✓  Strengths", Inches(0.55), Inches(1.67), Inches(5.8), Inches(0.34),
         font_size=Pt(13), bold=True, color=C_WHITE)
for i, pt in enumerate(strengths):
    add_text(s, f"• {pt}", Inches(0.55), Inches(2.1) + i * Inches(0.7),
             Inches(5.75), Inches(0.62), font_size=Pt(11), color=C_LIGHT_GRAY)

# Limitations card
add_rect(s, Inches(6.7), Inches(1.65), Inches(6.05), Inches(4.8), fill=C_CARD,
         line_color=C_YELLOW, line_width=Pt(1.5))
add_rect(s, Inches(6.7), Inches(1.65), Inches(6.05), Inches(0.38), fill=C_YELLOW)
add_text(s, "△  Limitations", Inches(6.85), Inches(1.67), Inches(5.8), Inches(0.34),
         font_size=Pt(13), bold=True, color=C_DARK_BG)
for i, pt in enumerate(limitations):
    add_text(s, f"• {pt}", Inches(6.85), Inches(2.1) + i * Inches(0.7),
             Inches(5.75), Inches(0.62), font_size=Pt(11), color=C_LIGHT_GRAY)

# Verdict
add_rect(s, Inches(0.4), Inches(6.55), Inches(12.4), Inches(0.85), fill=C_CARD)
add_rect(s, Inches(0.4), Inches(6.55), Inches(0.18), Inches(0.85), fill=C_ACCENT)
add_text(s,
    "Verdict: Claude Code is extremely powerful for prototyping and rapid development. "
    "The system functions correctly, but the output currently sits closer to a well-developed demo "
    "than a fully production-ready system. Human oversight for architecture and best practices remains essential.",
    Inches(0.72), Inches(6.6), Inches(12.0), Inches(0.75),
    font_size=Pt(12), color=C_LIGHT_GRAY, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Version 2 Roadmap
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "Version 2 — Planned Features")

add_text(s,
    "Version 2 features are intentionally designed to increase complexity and further test AI's ability "
    "to handle sophisticated software architecture.",
    Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.55),
    font_size=Pt(14), color=C_MID_GRAY, italic=True)

v2_features = [
    ("🔑 Single Sign-On",       "OAuth 2.0 / SAML integration with enterprise identity providers. Tests complex auth flows.",           C_ACCENT),
    ("📱 Mobile Responsive UI", "Full mobile layout overhaul. Tests whether AI can manage responsive breakpoints at scale.",             C_BLUE),
    ("🏷 Custom Task Fields",    "User-defined fields per task (dropdowns, dates, numbers). Tests dynamic schema management.",           C_GREEN),
    ("📊 Analytics Dashboard",  "Charts and reports for task throughput, team velocity, burndown. Tests data visualisation code.",       C_YELLOW),
    ("🔀 Sub-tasks",             "Hierarchical tasks with progress rollup to parent. Tests recursive data structures and UI.",            C_ORANGE),
]

for i, (title, desc, col) in enumerate(v2_features):
    y = Inches(1.85) + i * Inches(1.05)
    add_rect(s, Inches(0.4), y, Inches(12.4), Inches(0.92), fill=C_CARD,
             line_color=col, line_width=Pt(1))
    add_rect(s, Inches(0.4), y, Inches(0.18), Inches(0.92), fill=col)
    add_text(s, title, Inches(0.72), y + Inches(0.08), Inches(3.5), Inches(0.38),
             font_size=Pt(14), bold=True, color=col)
    add_text(s, desc, Inches(4.4), y + Inches(0.12), Inches(8.2), Inches(0.68),
             font_size=Pt(12), color=C_LIGHT_GRAY)
    add_rect(s, Inches(4.2), y + Inches(0.2), Inches(0.04), Inches(0.52),
             fill=RGBColor(0x4B, 0x55, 0x63))

add_text(s, "Goal: Determine whether AI can handle enterprise-grade architecture decisions",
         Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.3),
         font_size=Pt(12), color=C_ACCENT2, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Conclusion
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
add_rect(s, 0, Inches(1.5), SLIDE_W, Inches(4.5), fill=C_CARD)
add_rect(s, 0, Inches(1.5), Inches(0.3), Inches(4.5), fill=C_ACCENT)
accent_bar(s)

slide_title(s, "07 — Conclusion")

add_text(s, "Key Takeaways",
         Inches(0.6), Inches(1.65), Inches(12), Inches(0.5),
         font_size=Pt(22), bold=True, color=C_ACCENT2)

takeaways = [
    ("🚀 Speed",       "AI coding tools like Claude Code can build a full-stack web application in 4 days that would traditionally take weeks."),
    ("✅ Capability",  "Claude Code successfully implemented 16 features spanning auth, UI, APIs, email, real-time UX, and admin tooling."),
    ("⚠️ Oversight",   "The AI prioritises working code over best practices — human architectural direction is still required."),
    ("📋 Process",     "AI produces better results with clear prompts, structured requirements, and detailed documentation provided upfront."),
    ("🎯 Verdict",     "Version 1 is a well-developed demo. With more structured guidance, Version 2 aims to close the gap to production."),
]

for i, (icon_label, text) in enumerate(takeaways):
    y = Inches(2.2) + i * Inches(0.76)
    add_text(s, icon_label, Inches(0.6), y, Inches(1.1), Inches(0.6),
             font_size=Pt(13), bold=True, color=C_ACCENT2)
    add_text(s, text, Inches(1.8), y, Inches(10.8), Inches(0.6),
             font_size=Pt(13), color=C_LIGHT_GRAY)

add_rect(s, 0, Inches(6.4), SLIDE_W, Inches(1.1), fill=C_DARK_BG)
add_rect(s, 0, Inches(6.4), SLIDE_W, Inches(0.06), fill=C_ACCENT)
add_text(s, "Task Management System — Version 1",
         Inches(0.5), Inches(6.52), Inches(8), Inches(0.38),
         font_size=Pt(16), bold=True, color=C_WHITE)
add_text(s, "Built entirely with Claude Code  ·  March 2026",
         Inches(0.5), Inches(6.9), Inches(8), Inches(0.3),
         font_size=Pt(12), color=C_MID_GRAY)
add_text(s, "github.com/Zihern-sudo/TaskManagementSystem",
         Inches(8.5), Inches(6.7), Inches(4.3), Inches(0.4),
         font_size=Pt(12), color=C_ACCENT2, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════════════
out_path = "/home/user/TaskManagementSystem/TaskManagementSystem_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
